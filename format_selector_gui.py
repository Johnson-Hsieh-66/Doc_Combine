import os
import sys
import traceback
import logging
import glob
import argparse
from datetime import datetime
import tkinter as tk
from tkinter import ttk, messagebox
import io

# 設置標準輸出的編碼為UTF-8
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8')

# 設置日誌
def setup_logging():
    # 確保logs目錄存在
    logs_dir = "logs"
    if not os.path.exists(logs_dir):
        os.makedirs(logs_dir)
    
    # 創建日誌文件名（包含時間戳）
    log_file = os.path.join(logs_dir, f"format_selector_{datetime.now().strftime('%Y%m%d%H%M%S')}.log")
    
    # 配置日誌
    logging.basicConfig(
        level=logging.INFO,
        format='%(asctime)s - %(levelname)s - %(message)s',
        datefmt='%Y-%m-%d %H:%M:%S',
        handlers=[
            logging.FileHandler(log_file, encoding='utf-8'),
            logging.StreamHandler()  # 同時輸出到控制台
        ]
    )
    
    logger = logging.getLogger('格式選擇器')
    logger.info(f"日誌文件已創建: {os.path.abspath(log_file)}")
    return logger

# 創建日誌記錄器
logger = setup_logging()

# 從merge_files.py整合的函數
def generate_ppt(output_file, status_callback=None):
    """生成PPT文件並保存到指定路徑"""
    # 確保輸出目錄存在
    output_dir = os.path.dirname(output_file)
    if output_dir and not os.path.exists(output_dir):
        os.makedirs(output_dir)
    
    try:
        # 嘗試使用高級方法（需要 Windows 和 PowerPoint）
        import win32com.client
        
        # 獲取 docs 文件夾中的所有 PPT 文件
        ppt_files = glob.glob("docs/*.ppt*")
        
        if not ppt_files:
            msg = "docs 文件夾中沒有找到 PPT 文件"
            if status_callback:
                status_callback(msg)
            logger.warning(msg)
            return False
        
        # 啟動 PowerPoint 應用程序
        ppt_app = win32com.client.Dispatch("PowerPoint.Application")
        
        # 使用第一個文件作為基礎，而不是創建新的空演示文稿
        msg = f"正在使用第一個文件作為基礎: {ppt_files[0]}"
        if status_callback:
            status_callback(msg)
        logger.info(msg)
        
        base_path = os.path.abspath(ppt_files[0])
        merged_presentation = ppt_app.Presentations.Open(base_path)
        
        # 遍歷剩餘的 PPT 文件並合併
        for ppt_file in ppt_files[1:]:
            msg = f"正在處理: {ppt_file}"
            if status_callback:
                status_callback(msg)
            logger.info(msg)
            
            try:
                # 打開當前 PPT 文件
                abs_path = os.path.abspath(ppt_file)
                current_presentation = ppt_app.Presentations.Open(abs_path)
                
                # 獲取幻燈片數量
                slide_count = current_presentation.Slides.Count
                
                # 複製所有幻燈片到合併的演示文稿
                for i in range(1, slide_count + 1):
                    current_presentation.Slides(i).Copy()
                    merged_presentation.Slides.Paste()
                
                # 關閉當前演示文稿而不保存
                current_presentation.Close()
            except Exception as e:
                msg = f"處理文件 {ppt_file} 時出錯: {e}"
                if status_callback:
                    status_callback(msg)
                logger.error(msg)
        
        # 保存合併後的 PPT 到指定路徑
        output_abs_path = os.path.abspath(output_file)
        merged_presentation.SaveAs(output_abs_path)
        merged_presentation.Close()
        
        # 關閉 PowerPoint 應用程序
        ppt_app.Quit()
        msg = f"已成功合併所有 PPT 文件到: {output_file}"
        if status_callback:
            status_callback(msg)
        logger.info(msg)
        return True
        
    except ImportError as e:
        msg = f"導入錯誤: {e}"
        if status_callback:
            status_callback(msg)
        logger.error(msg)
        
        # 如果無法導入 win32com，則使用 python-pptx（功能有限）
        try:
            from pptx import Presentation
            
            # 獲取 docs 文件夾中的所有 PPT 文件
            ppt_files = glob.glob("docs/*.ppt*")
            
            if not ppt_files:
                msg = "docs 文件夾中沒有找到 PPT 文件"
                if status_callback:
                    status_callback(msg)
                logger.warning(msg)
                return False
            
            # 創建一個新的演示文稿
            merged_ppt = Presentation()
            
            # 遍歷所有 PPT 文件並合併
            for ppt_file in ppt_files:
                msg = f"正在處理: {ppt_file}"
                if status_callback:
                    status_callback(msg)
                logger.info(msg)
                
                try:
                    # 打開當前 PPT 文件
                    current_ppt = Presentation(ppt_file)
                    
                    # 複製每一張幻燈片到新的演示文稿
                    for slide in current_ppt.slides:
                        # 複製幻燈片布局
                        slide_layout = merged_ppt.slide_layouts[0]  # 使用默認布局
                        new_slide = merged_ppt.slides.add_slide(slide_layout)
                        
                        # 複製所有形狀
                        for shape in slide.shapes:
                            # 這裡我們只能複製基本元素，複雜元素可能需要更詳細的處理
                            if shape.has_text_frame:
                                for paragraph in shape.text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        text_shape = new_slide.shapes.add_textbox(
                                            shape.left, shape.top, shape.width, shape.height
                                        )
                                        text_frame = text_shape.text_frame
                                        p = text_frame.add_paragraph()
                                        p.text = run.text
                                        # 注意：這裡沒有複製格式，如果需要可以添加更多代碼
                except Exception as e:
                    msg = f"處理文件 {ppt_file} 時出錯: {e}"
                    if status_callback:
                        status_callback(msg)
                    logger.error(msg)
            
            # 保存合併後的 PPT 到指定路徑
            merged_ppt.save(output_file)
            msg = f"已成功合併所有 PPT 文件到: {output_file}"
            if status_callback:
                status_callback(msg)
            logger.info(msg)
            
            msg = "注意：使用 python-pptx 合併可能會丟失一些格式和效果"
            if status_callback:
                status_callback(msg)
            logger.warning(msg)
            return True
            
        except ImportError:
            msg = "錯誤：需要安裝 python-pptx 庫才能合併 PPT 文件"
            if status_callback:
                status_callback(msg)
            logger.error(msg)
            
            msg = "請運行: pip install python-pptx"
            if status_callback:
                status_callback(msg)
            logger.info(msg)
            return False

def generate_pdf(output_file, status_callback=None):
    """生成PDF文件並保存到指定路徑"""
    # 確保輸出目錄存在
    output_dir = os.path.dirname(output_file)
    if output_dir and not os.path.exists(output_dir):
        os.makedirs(output_dir)
    
    try:
        from PyPDF2 import PdfMerger
        
        # 獲取 docs 文件夾中的所有 PDF 文件
        pdf_files = glob.glob("docs/*.pdf")
        
        if not pdf_files:
            msg = "docs 文件夾中沒有找到 PDF 文件"
            if status_callback:
                status_callback(msg)
            logger.warning(msg)
            return False
        
        # 創建 PDF 合併器
        merger = PdfMerger()
        
        # 遍歷所有 PDF 文件並合併
        for pdf_file in pdf_files:
            msg = f"正在處理: {pdf_file}"
            if status_callback:
                status_callback(msg)
            logger.info(msg)
            
            try:
                merger.append(pdf_file)
            except Exception as e:
                msg = f"處理文件 {pdf_file} 時出錯: {e}"
                if status_callback:
                    status_callback(msg)
                logger.error(msg)
        
        # 保存合併後的 PDF 到指定路徑
        merger.write(output_file)
        merger.close()
        msg = f"已成功合併所有 PDF 文件到: {output_file}"
        if status_callback:
            status_callback(msg)
        logger.info(msg)
        return True
        
    except ImportError:
        msg = "錯誤：需要安裝 PyPDF2 庫才能合併 PDF 文件"
        if status_callback:
            status_callback(msg)
        logger.error(msg)
        
        msg = "請運行: pip install PyPDF2"
        if status_callback:
            status_callback(msg)
        logger.info(msg)
        return False

class FormatSelectorApp:
    def __init__(self, root):
        self.root = root
        self.root.title("文件格式選擇器")
        self.root.geometry("400x350")
        self.root.resizable(False, False)
        
        # 設置格式變量
        self.format_var = tk.StringVar(value="ppt")
        
        # 創建界面元素
        self.create_widgets()
        
        # 居中窗口
        self.center_window()
        
        logger.info("Tkinter界面已初始化")
    
    def create_widgets(self):
        # 標題標籤
        title_label = ttk.Label(
            self.root, 
            text="請選擇要生成的文件格式",
            font=("Microsoft Sans Serif", 12, "bold")
        )
        title_label.pack(pady=(20, 15))
        
        # 格式選擇框架
        format_frame = ttk.LabelFrame(self.root, text="輸出格式")
        format_frame.pack(fill="x", padx=50, pady=10)
        
        # PPT單選按鈕
        self.ppt_radio = ttk.Radiobutton(
            format_frame,
            text="PowerPoint (.pptx)",
            variable=self.format_var,
            value="ppt",
            command=self.on_format_changed
        )
        self.ppt_radio.pack(anchor="w", padx=20, pady=(10, 5))
        
        # PDF單選按鈕
        self.pdf_radio = ttk.Radiobutton(
            format_frame,
            text="PDF (.pdf)",
            variable=self.format_var,
            value="pdf",
            command=self.on_format_changed
        )
        self.pdf_radio.pack(anchor="w", padx=20, pady=(5, 10))
        
        # 狀態文本框
        self.status_text = tk.Text(
            self.root,
            height=5,
            width=40,
            wrap="word",
            state="normal"
        )
        self.status_text.pack(padx=50, pady=15, fill="x")
        
        # 添加滾動條
        scrollbar = ttk.Scrollbar(self.status_text)
        scrollbar.pack(side="right", fill="y")
        self.status_text.config(yscrollcommand=scrollbar.set)
        scrollbar.config(command=self.status_text.yview)
        
        # 生成按鈕
        self.generate_button = ttk.Button(
            self.root,
            text="生成文件",
            command=self.on_generate_click
        )
        self.generate_button.pack(pady=15)
        
        # 初始化狀態信息
        self.update_status("準備就緒。請選擇格式並點擊「生成文件」按鈕。")
    
    def center_window(self):
        self.root.update_idletasks()
        width = self.root.winfo_width()
        height = self.root.winfo_height()
        x = (self.root.winfo_screenwidth() // 2) - (width // 2)
        y = (self.root.winfo_screenheight() // 2) - (height // 2)
        self.root.geometry(f"{width}x{height}+{x}+{y}")
    
    def on_format_changed(self):
        format_type = self.format_var.get()
        logger.info(f"用戶選擇了格式: {format_type}")
    
    def update_status(self, message):
        self.status_text.config(state="normal")
        self.status_text.insert("end", f"{message}\n")
        self.status_text.see("end")
        self.status_text.config(state="disabled")
        # 更新界面
        self.root.update()
    
    def on_generate_click(self):
        try:
            format_type = self.format_var.get()
            self.update_status(f"正在生成{format_type.upper()}文件，請稍候...")
            logger.info(f"開始生成{format_type.upper()}文件")
            
            # 禁用生成按鈕
            self.generate_button.config(state="disabled")
            self.root.update()
            
            # 獲取腳本路徑
            if getattr(sys, 'frozen', False):
                # 如果是打包後的可執行文件
                script_dir = os.path.dirname(sys.executable)
            else:
                # 如果是Python腳本
                script_dir = os.path.dirname(os.path.abspath(__file__))
            
            logger.info(f"腳本目錄: {script_dir}")
            
            # 檢查docs目錄是否存在
            docs_dir = os.path.join(script_dir, "docs")
            if not os.path.exists(docs_dir):
                os.makedirs(docs_dir)
                logger.warning(f"docs目錄不存在，已創建: {docs_dir}")
                self.update_status(f"警告: docs目錄不存在，已創建空目錄。")
                self.update_status(f"請將要合併的文件放入此目錄: {docs_dir}")
            
            # 確保output目錄存在
            output_dir = os.path.join(script_dir, "output")
            if not os.path.exists(output_dir):
                os.makedirs(output_dir)
                logger.info(f"output目錄不存在，已創建: {output_dir}")
            
            # 構建輸出文件路徑
            timestamp = datetime.now().strftime("%Y%m%d%H%M%S")
            output_file = os.path.join(output_dir, f"output_{timestamp}.{format_type}")
            logger.info(f"輸出文件路徑: {output_file}")
            
            # 切換到腳本目錄，確保相對路徑正確
            os.chdir(script_dir)
            
            # 直接調用合併函數，而不是使用subprocess
            success = False
            if format_type == "ppt":
                success = generate_ppt(output_file, self.update_status)
            elif format_type == "pdf":
                success = generate_pdf(output_file, self.update_status)
            
            # 啟用生成按鈕
            self.generate_button.config(state="normal")
            
            # 檢查執行結果
            if success:
                logger.info(f"文件生成成功: {output_file}")
                self.update_status(f"{format_type.upper()}文件已生成成功！")
                self.update_status(f"文件位置：{output_file}")
                
                # 詢問是否打開文件所在的文件夾
                answer = messagebox.askyesno(
                    "生成成功",
                    f"{format_type.upper()}文件已生成成功！\n\n文件位置：{output_file}\n\n是否打開文件所在的文件夾？"
                )
                
                if answer:
                    try:
                        os.startfile(output_dir)
                        logger.info(f"已打開文件夾: {output_dir}")
                    except Exception as e:
                        logger.error(f"打開文件夾失敗: {str(e)}")
                        self.update_status(f"打開文件夾時發生錯誤: {str(e)}")
                        self.update_status("請手動瀏覽到以下位置查看文件:")
                        self.update_status(output_dir)
            else:
                logger.error(f"生成失敗")
                self.update_status(f"生成失敗")
                self.update_status("請檢查logs目錄下的日誌文件以獲取詳細錯誤信息")
                
                messagebox.showerror(
                    "生成失敗",
                    f"生成{format_type.upper()}文件失敗！\n\n請檢查logs目錄下的日誌文件以獲取詳細信息。"
                )
                
        except Exception as e:
            error_details = traceback.format_exc()
            logger.error(f"發生異常: {str(e)}\n{error_details}")
            self.update_status(f"發生錯誤: {str(e)}")
            self.update_status("詳細錯誤信息已寫入日誌文件，請查看logs目錄")
            
            # 啟用生成按鈕
            self.generate_button.config(state="normal")
            
            messagebox.showerror(
                "錯誤",
                f"處理過程中發生錯誤！\n\n錯誤信息：{str(e)}\n\n詳細錯誤信息已寫入日誌文件，請查看logs目錄。"
            )

def main():
    try:
        root = tk.Tk()
        app = FormatSelectorApp(root)
        root.mainloop()
    except Exception as e:
        error_details = traceback.format_exc()
        logger.error(f"程序啟動時發生嚴重錯誤: {str(e)}\n{error_details}")
        
        # 嘗試顯示錯誤對話框
        try:
            messagebox.showerror(
                "嚴重錯誤",
                f"程序啟動時發生嚴重錯誤！\n\n錯誤信息：{str(e)}\n\n詳細錯誤信息已寫入日誌文件。"
            )
        except:
            # 如果無法顯示對話框，則使用控制台輸出
            print(f"\n程序啟動時發生嚴重錯誤: {str(e)}")
            print("\n詳細錯誤信息:")
            print(error_details)
            input("\n按Enter鍵退出...")

if __name__ == "__main__":
    main() 