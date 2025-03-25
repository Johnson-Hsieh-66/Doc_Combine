import os
import glob
import argparse
from datetime import datetime
import shutil
import sys
import io

# 設置標準輸出的編碼為UTF-8
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8')

def merge_ppt_files():
    # 創建時間戳作為文件名
    timestamp = datetime.now().strftime("%Y%m%d%H%M%S")
    output_filename = f"output/{timestamp}.pptx"
    
    # 確保 output 文件夾存在
    if not os.path.exists("output"):
        os.makedirs("output")
    
    try:
        # 嘗試使用高級方法（需要 Windows 和 PowerPoint）
        import win32com.client
        
        # 獲取 docs 文件夾中的所有 PPT 文件
        ppt_files = glob.glob("docs/*.ppt*")
        
        if not ppt_files:
            print("docs 文件夾中沒有找到 PPT 文件")
            return
        
        # 啟動 PowerPoint 應用程序
        ppt_app = win32com.client.Dispatch("PowerPoint.Application")
        
        # 使用第一個文件作為基礎，而不是創建新的空演示文稿
        print(f"正在使用第一個文件作為基礎: {ppt_files[0]}")
        base_path = os.path.abspath(ppt_files[0])
        merged_presentation = ppt_app.Presentations.Open(base_path)
        
        # 遍歷剩餘的 PPT 文件並合併
        for ppt_file in ppt_files[1:]:
            print(f"正在處理: {ppt_file}")
            try:
                # 打開當前 PPT 文件
                abs_path = os.path.abspath(ppt_file)
                current_presentation = ppt_app.Presentations.Open(abs_path)
                
                # 獲取幻燈片數量
                slide_count = current_presentation.Slides.Count
                
                # 複製所有幻燈片到合併的演示文稿
                for i in range(1, slide_count + 1):
                    current_presentation.Slides(i).Copy()
                    merged_presentation.Slides.Paste()
                
                # 關閉當前演示文稿而不保存
                current_presentation.Close()
            except Exception as e:
                print(f"處理文件 {ppt_file} 時出錯: {e}")
        
        # 保存合併後的 PPT
        merged_presentation.SaveAs(os.path.abspath(output_filename))
        merged_presentation.Close()
        
        # 關閉 PowerPoint 應用程序
        ppt_app.Quit()
        print(f"已成功合併所有 PPT 文件到: {output_filename}")
        
    except ImportError as e:
        print(f"導入錯誤: {e}")
        # 如果無法導入 win32com，則使用 python-pptx（功能有限）
        try:
            from pptx import Presentation
            
            # 獲取 docs 文件夾中的所有 PPT 文件
            ppt_files = glob.glob("docs/*.ppt*")
            
            if not ppt_files:
                print("docs 文件夾中沒有找到 PPT 文件")
                return
            
            # 創建一個新的演示文稿
            merged_ppt = Presentation()
            
            # 遍歷所有 PPT 文件並合併
            for ppt_file in ppt_files:
                print(f"正在處理: {ppt_file}")
                try:
                    # 打開當前 PPT 文件
                    current_ppt = Presentation(ppt_file)
                    
                    # 複製每一張幻燈片到新的演示文稿
                    for slide in current_ppt.slides:
                        # 複製幻燈片布局
                        slide_layout = merged_ppt.slide_layouts[0]  # 使用默認布局
                        new_slide = merged_ppt.slides.add_slide(slide_layout)
                        
                        # 複製所有形狀
                        for shape in slide.shapes:
                            # 這裡我們只能複製基本元素，複雜元素可能需要更詳細的處理
                            if shape.has_text_frame:
                                for paragraph in shape.text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        text_shape = new_slide.shapes.add_textbox(
                                            shape.left, shape.top, shape.width, shape.height
                                        )
                                        text_frame = text_shape.text_frame
                                        p = text_frame.add_paragraph()
                                        p.text = run.text
                                        # 注意：這裡沒有複製格式，如果需要可以添加更多代碼
                except Exception as e:
                    print(f"處理文件 {ppt_file} 時出錯: {e}")
            
            # 保存合併後的 PPT
            merged_ppt.save(output_filename)
            print(f"已成功合併所有 PPT 文件到: {output_filename}")
            print("注意：使用 python-pptx 合併可能會丟失一些格式和效果")
            
        except ImportError:
            print("錯誤：需要安裝 python-pptx 庫才能合併 PPT 文件")
            print("請運行: pip install python-pptx")

def merge_pdf_files():
    # 創建時間戳作為文件名
    timestamp = datetime.now().strftime("%Y%m%d%H%M%S")
    output_filename = f"output/{timestamp}.pdf"
    
    # 確保 output 文件夾存在
    if not os.path.exists("output"):
        os.makedirs("output")
    
    try:
        from PyPDF2 import PdfMerger
        
        # 獲取 docs 文件夾中的所有 PDF 文件
        pdf_files = glob.glob("docs/*.pdf")
        
        if not pdf_files:
            print("docs 文件夾中沒有找到 PDF 文件")
            return
        
        # 創建 PDF 合併器
        merger = PdfMerger()
        
        # 遍歷所有 PDF 文件並合併
        for pdf_file in pdf_files:
            print(f"正在處理: {pdf_file}")
            try:
                merger.append(pdf_file)
            except Exception as e:
                print(f"處理文件 {pdf_file} 時出錯: {e}")
        
        # 保存合併後的 PDF
        merger.write(output_filename)
        merger.close()
        print(f"已成功合併所有 PDF 文件到: {output_filename}")
        
    except ImportError:
        print("錯誤：需要安裝 PyPDF2 庫才能合併 PDF 文件")
        print("請運行: pip install PyPDF2")

def main():
    # 添加命令行參數支持
    parser = argparse.ArgumentParser(description='生成PPT或PDF文件')
    parser.add_argument('--format', type=str, default='ppt', help='輸出格式 (ppt 或 pdf)')
    parser.add_argument('--output', type=str, default='output.ppt', help='輸出文件名')
    
    args = parser.parse_args()
    
    # 根據格式調用不同的處理函數
    if args.format.lower() == 'ppt':
        generate_ppt(args.output)
    elif args.format.lower() == 'pdf':
        generate_pdf(args.output)
    else:
        print(f"不支持的格式: {args.format}")
        return 1
    
    print(f"文件已生成: {args.output}")
    return 0

def generate_ppt(output_file):
    """生成PPT文件並保存到指定路徑"""
    # 確保輸出目錄存在
    output_dir = os.path.dirname(output_file)
    if output_dir and not os.path.exists(output_dir):
        os.makedirs(output_dir)
    
    try:
        # 嘗試使用高級方法（需要 Windows 和 PowerPoint）
        import win32com.client
        
        # 獲取 docs 文件夾中的所有 PPT 文件
        ppt_files = glob.glob("docs/*.ppt*")
        
        if not ppt_files:
            print("docs 文件夾中沒有找到 PPT 文件")
            return
        
        # 啟動 PowerPoint 應用程序
        ppt_app = win32com.client.Dispatch("PowerPoint.Application")
        
        # 使用第一個文件作為基礎，而不是創建新的空演示文稿
        print(f"正在使用第一個文件作為基礎: {ppt_files[0]}")
        base_path = os.path.abspath(ppt_files[0])
        merged_presentation = ppt_app.Presentations.Open(base_path)
        
        # 遍歷剩餘的 PPT 文件並合併
        for ppt_file in ppt_files[1:]:
            print(f"正在處理: {ppt_file}")
            try:
                # 打開當前 PPT 文件
                abs_path = os.path.abspath(ppt_file)
                current_presentation = ppt_app.Presentations.Open(abs_path)
                
                # 獲取幻燈片數量
                slide_count = current_presentation.Slides.Count
                
                # 複製所有幻燈片到合併的演示文稿
                for i in range(1, slide_count + 1):
                    current_presentation.Slides(i).Copy()
                    merged_presentation.Slides.Paste()
                
                # 關閉當前演示文稿而不保存
                current_presentation.Close()
            except Exception as e:
                print(f"處理文件 {ppt_file} 時出錯: {e}")
        
        # 保存合併後的 PPT 到指定路徑
        output_abs_path = os.path.abspath(output_file)
        merged_presentation.SaveAs(output_abs_path)
        merged_presentation.Close()
        
        # 關閉 PowerPoint 應用程序
        ppt_app.Quit()
        print(f"已成功合併所有 PPT 文件到: {output_file}")
        
    except ImportError as e:
        print(f"導入錯誤: {e}")
        # 如果無法導入 win32com，則使用 python-pptx（功能有限）
        try:
            from pptx import Presentation
            
            # 獲取 docs 文件夾中的所有 PPT 文件
            ppt_files = glob.glob("docs/*.ppt*")
            
            if not ppt_files:
                print("docs 文件夾中沒有找到 PPT 文件")
                return
            
            # 創建一個新的演示文稿
            merged_ppt = Presentation()
            
            # 遍歷所有 PPT 文件並合併
            for ppt_file in ppt_files:
                print(f"正在處理: {ppt_file}")
                try:
                    # 打開當前 PPT 文件
                    current_ppt = Presentation(ppt_file)
                    
                    # 複製每一張幻燈片到新的演示文稿
                    for slide in current_ppt.slides:
                        # 複製幻燈片布局
                        slide_layout = merged_ppt.slide_layouts[0]  # 使用默認布局
                        new_slide = merged_ppt.slides.add_slide(slide_layout)
                        
                        # 複製所有形狀
                        for shape in slide.shapes:
                            # 這裡我們只能複製基本元素，複雜元素可能需要更詳細的處理
                            if shape.has_text_frame:
                                for paragraph in shape.text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        text_shape = new_slide.shapes.add_textbox(
                                            shape.left, shape.top, shape.width, shape.height
                                        )
                                        text_frame = text_shape.text_frame
                                        p = text_frame.add_paragraph()
                                        p.text = run.text
                                        # 注意：這裡沒有複製格式，如果需要可以添加更多代碼
                except Exception as e:
                    print(f"處理文件 {ppt_file} 時出錯: {e}")
            
            # 保存合併後的 PPT 到指定路徑
            merged_ppt.save(output_file)
            print(f"已成功合併所有 PPT 文件到: {output_file}")
            print("注意：使用 python-pptx 合併可能會丟失一些格式和效果")
            
        except ImportError:
            print("錯誤：需要安裝 python-pptx 庫才能合併 PPT 文件")
            print("請運行: pip install python-pptx")

def generate_pdf(output_file):
    """生成PDF文件並保存到指定路徑"""
    # 確保輸出目錄存在
    output_dir = os.path.dirname(output_file)
    if output_dir and not os.path.exists(output_dir):
        os.makedirs(output_dir)
    
    try:
        from PyPDF2 import PdfMerger
        
        # 獲取 docs 文件夾中的所有 PDF 文件
        pdf_files = glob.glob("docs/*.pdf")
        
        if not pdf_files:
            print("docs 文件夾中沒有找到 PDF 文件")
            return
        
        # 創建 PDF 合併器
        merger = PdfMerger()
        
        # 遍歷所有 PDF 文件並合併
        for pdf_file in pdf_files:
            print(f"正在處理: {pdf_file}")
            try:
                merger.append(pdf_file)
            except Exception as e:
                print(f"處理文件 {pdf_file} 時出錯: {e}")
        
        # 保存合併後的 PDF 到指定路徑
        merger.write(output_file)
        merger.close()
        print(f"已成功合併所有 PDF 文件到: {output_file}")
        
    except ImportError:
        print("錯誤：需要安裝 PyPDF2 庫才能合併 PDF 文件")
        print("請運行: pip install PyPDF2")

if __name__ == "__main__":
    sys.exit(main()) 